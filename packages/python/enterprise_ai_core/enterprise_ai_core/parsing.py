from __future__ import annotations

from email import policy
from email.parser import BytesParser
from io import BytesIO
import re
from uuid import uuid4

import fitz
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from PIL import Image
from pypdf import PdfReader
import pytesseract
from pytesseract import Output, TesseractNotFoundError

from .config import get_settings
from .schemas import CanonicalBlock, CanonicalDocument, RuntimeHealthResponse


def detect_language(text: str) -> str:
    vietnamese_markers = (
        "\u0103\u00e2\u0111\u00ea\u00f4\u01a1\u01b0\u00e1\u00e0\u1ea3\u00e3\u1ea1"
        "\u1ebf\u1ec1\u1ec3\u1ec5\u1ec7\u1ed1\u1ed3\u1ed5\u1ed7\u1ed9"
    )
    if any(character in text.lower() for character in vietnamese_markers):
        return "vi"
    if re.search(r"[A-Za-z]", text):
        return "en"
    return "unknown"


def configure_tesseract() -> None:
    settings = get_settings()
    if settings.tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd


def ocr_runtime_health(service_name: str) -> RuntimeHealthResponse:
    settings = get_settings()
    try:
        configure_tesseract()
        version = str(pytesseract.get_tesseract_version())
        return RuntimeHealthResponse(
            service=service_name,
            runtime="ocr",
            status="ok",
            detail="Tesseract is available",
            metadata={
                "engine": settings.ocr_engine,
                "languages": settings.ocr_languages,
                "version": version,
            },
        )
    except Exception as exc:
        return RuntimeHealthResponse(
            service=service_name,
            runtime="ocr",
            status="error",
            detail=str(exc),
            metadata={"engine": settings.ocr_engine, "languages": settings.ocr_languages},
        )


def extract_text_blocks(
    file_name: str,
    content_type: str,
    payload: bytes,
) -> tuple[list[CanonicalBlock], bool, bool, dict]:
    settings = get_settings()
    suffix = file_name.lower().rsplit(".", 1)[-1] if "." in file_name else ""
    blocks: list[CanonicalBlock] = []
    warnings: list[str] = []
    ocr_required = False
    ocr_applied = False
    parse_metadata: dict[str, object] = {"format": suffix or content_type}

    try:
        if content_type == "text/plain" or suffix in {"txt", "log"}:
            text = payload.decode("utf-8", errors="ignore")
            blocks = text_to_blocks(text, "paragraph")
        elif suffix in {"md", "markdown"} or content_type in {
            "text/markdown",
            "text/x-markdown",
        }:
            text = payload.decode("utf-8", errors="ignore")
            blocks = markdown_to_blocks(text)
        elif suffix in {"html", "htm"} or content_type == "text/html":
            blocks = html_to_blocks(payload.decode("utf-8", errors="ignore"))
        elif suffix == "pdf" or content_type == "application/pdf":
            blocks = pdf_to_blocks(payload)
            extracted_chars = sum(len(item.text.strip()) for item in blocks)
            ocr_required = extracted_chars < settings.ocr_min_characters
            if ocr_required and settings.ocr_engine == "tesseract":
                ocr_blocks = pdf_to_ocr_blocks(payload)
                if ocr_blocks:
                    blocks = ocr_blocks
                    ocr_applied = True
                else:
                    warnings.append("pdf_ocr_fallback_returned_no_text")
        elif suffix == "docx" or (
            content_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            blocks = docx_to_blocks(payload)
        elif suffix == "xlsx" or (
            content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ):
            blocks = xlsx_to_blocks(payload)
        elif suffix == "pptx" or (
            content_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            blocks = pptx_to_blocks(payload)
        elif suffix == "eml" or content_type == "message/rfc822":
            blocks = eml_to_blocks(payload)
        elif suffix in {"png", "jpg", "jpeg", "tif", "tiff", "bmp", "webp"} or content_type.startswith(
            "image/"
        ):
            ocr_required = True
            if settings.ocr_engine == "tesseract":
                blocks = image_to_blocks(payload)
                ocr_applied = bool(blocks)
        else:
            warnings.append(f"unsupported_format_fallback:{suffix or content_type}")
            text = payload.decode("utf-8", errors="ignore")
            blocks = text_to_blocks(text, "paragraph")
    except Exception as exc:
        warnings.append(f"parser_error:{type(exc).__name__}")
        text = payload.decode("utf-8", errors="ignore")
        blocks = text_to_blocks(text, "paragraph")

    finalize_offsets(blocks)
    parse_quality_score = calculate_parse_quality(
        blocks=blocks,
        warnings=warnings,
        ocr_required=ocr_required,
        ocr_applied=ocr_applied,
    )
    parse_metadata.update(
        {
            "warning_count": len(warnings),
            "warnings": warnings,
            "block_count": len(blocks),
            "parse_quality_score": parse_quality_score,
        }
    )
    return blocks, ocr_required, ocr_applied, parse_metadata


def html_to_text(html: str) -> str:
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text("\n")


def html_to_blocks(html: str) -> list[CanonicalBlock]:
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()

    blocks: list[CanonicalBlock] = []
    heading_path: list[str] = []
    order_index = 0
    for element in soup.find_all(
        ["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "blockquote", "tr"]
    ):
        text = " ".join(element.get_text(" ", strip=True).split())
        if not text:
            continue
        if element.name and element.name.startswith("h"):
            level = int(element.name[1])
            heading_path = heading_path[: max(0, level - 1)]
            heading_path.append(text)
            blocks.append(
                make_block(
                    order_index=order_index,
                    block_type="heading",
                    text=text,
                    heading=text,
                    heading_path=list(heading_path),
                )
            )
        elif element.name == "tr":
            cells = [cell.get_text(" ", strip=True) for cell in element.find_all(["td", "th"])]
            if not any(cells):
                continue
            blocks.append(
                make_block(
                    order_index=order_index,
                    block_type="table-row",
                    text=" | ".join(cells),
                    heading=heading_path[-1] if heading_path else None,
                    heading_path=list(heading_path),
                    metadata={"table_columns": cells},
                )
            )
        else:
            blocks.append(
                make_block(
                    order_index=order_index,
                    block_type="paragraph",
                    text=text,
                    heading=heading_path[-1] if heading_path else None,
                    heading_path=list(heading_path),
                )
            )
        order_index += 1
    if not blocks:
        return text_to_blocks(html_to_text(html), "paragraph")
    return blocks


def text_to_blocks(text: str, block_type: str) -> list[CanonicalBlock]:
    blocks: list[CanonicalBlock] = []
    lines = [line.strip() for line in text.splitlines()]
    order_index = 0
    for line in lines:
        if not line:
            continue
        blocks.append(
            make_block(
                order_index=order_index,
                block_type=block_type,
                text=line,
            )
        )
        order_index += 1
    return blocks


def markdown_to_blocks(text: str) -> list[CanonicalBlock]:
    blocks: list[CanonicalBlock] = []
    order_index = 0
    heading_path: list[str] = []
    current_heading: str | None = None
    current_body: list[str] = []

    def flush_current() -> None:
        nonlocal order_index, current_body
        body = "\n".join(item for item in current_body if item).strip()
        if body:
            blocks.append(
                make_block(
                    order_index=order_index,
                    block_type="section",
                    text=body,
                    heading=current_heading,
                    heading_path=list(heading_path),
                )
            )
            order_index += 1
        current_body = []

    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("#"):
            flush_current()
            level = len(line) - len(line.lstrip("#"))
            heading_text = line.lstrip("#").strip()
            heading_path[:] = heading_path[: max(0, level - 1)]
            heading_path.append(heading_text)
            current_heading = heading_text
            blocks.append(
                make_block(
                    order_index=order_index,
                    block_type="heading",
                    text=heading_text,
                    heading=heading_text,
                    heading_path=list(heading_path),
                )
            )
            order_index += 1
        else:
            current_body.append(line)

    flush_current()
    if not blocks:
        return text_to_blocks(text, "paragraph")
    return blocks


def pdf_to_blocks(payload: bytes) -> list[CanonicalBlock]:
    blocks: list[CanonicalBlock] = []
    reader = PdfReader(BytesIO(payload))
    for page_index, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        blocks.append(
            make_block(
                order_index=page_index,
                block_type="page",
                text=text,
                heading=f"Page {page_index + 1}",
                heading_path=[f"Page {page_index + 1}"],
                page_start=page_index + 1,
                page_end=page_index + 1,
                parse_quality_score=0.95,
            )
        )
    return blocks


def pdf_to_ocr_blocks(payload: bytes) -> list[CanonicalBlock]:
    settings = get_settings()
    blocks: list[CanonicalBlock] = []
    document = fitz.open(stream=payload, filetype="pdf")
    scale = settings.ocr_render_dpi / 72
    matrix = fitz.Matrix(scale, scale)

    for page_index, page in enumerate(document):
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        image = Image.open(BytesIO(pixmap.tobytes("png")))
        text, confidence = ocr_image(image)
        if not text.strip():
            continue
        blocks.append(
            make_block(
                order_index=page_index,
                block_type="ocr-page",
                text=text,
                heading=f"Page {page_index + 1}",
                heading_path=[f"Page {page_index + 1}"],
                page_start=page_index + 1,
                page_end=page_index + 1,
                parse_quality_score=confidence / 100 if confidence else 0.55,
                ocr_confidence=confidence,
                metadata={"ocr_confidence": confidence},
            )
        )
    return blocks


def image_to_blocks(payload: bytes) -> list[CanonicalBlock]:
    image = Image.open(BytesIO(payload))
    text, confidence = ocr_image(image)
    if not text.strip():
        return []
    return [
        make_block(
            order_index=0,
            block_type="ocr-image",
            text=text,
            parse_quality_score=confidence / 100 if confidence else 0.55,
            ocr_confidence=confidence,
            metadata={"ocr_confidence": confidence},
        )
    ]


def ocr_image(image: Image.Image) -> tuple[str, float]:
    settings = get_settings()
    configure_tesseract()
    try:
        data = pytesseract.image_to_data(
            image,
            lang=settings.ocr_languages,
            output_type=Output.DICT,
        )
    except TesseractNotFoundError:
        return "", 0.0

    tokens: list[str] = []
    confidences: list[float] = []
    for raw_text, raw_confidence in zip(data.get("text", []), data.get("conf", []), strict=False):
        text = (raw_text or "").strip()
        if not text:
            continue
        tokens.append(text)
        try:
            confidence = float(raw_confidence)
        except (TypeError, ValueError):
            confidence = -1
        if confidence >= 0:
            confidences.append(confidence)

    joined_text = " ".join(tokens).strip()
    average_confidence = sum(confidences) / len(confidences) if confidences else 0.0
    return joined_text, average_confidence


def docx_to_blocks(payload: bytes) -> list[CanonicalBlock]:
    blocks: list[CanonicalBlock] = []
    document = DocxDocument(BytesIO(payload))
    order_index = 0
    heading_path: list[str] = []

    for paragraph in document.paragraphs:
        text = " ".join(paragraph.text.strip().split())
        if not text:
            continue
        style_name = getattr(paragraph.style, "name", "") or ""
        if "Heading" in style_name:
            level = parse_heading_level(style_name)
            heading_path = heading_path[: max(0, level - 1)]
            heading_path.append(text)
            block_type = "heading"
        else:
            block_type = "paragraph"
        blocks.append(
            make_block(
                order_index=order_index,
                block_type=block_type,
                text=text,
                heading=heading_path[-1] if heading_path else None,
                heading_path=list(heading_path),
            )
        )
        order_index += 1

    for table_index, table in enumerate(document.tables):
        for row_index, row in enumerate(table.rows):
            cells = [" ".join(cell.text.split()) for cell in row.cells]
            if not any(cells):
                continue
            blocks.append(
                make_block(
                    order_index=order_index,
                    block_type="table-row",
                    text=" | ".join(cells),
                    heading=heading_path[-1] if heading_path else None,
                    heading_path=list(heading_path),
                    table_id=f"docx-table-{table_index}",
                    row_index=row_index,
                    metadata={"table_columns": cells},
                )
            )
            order_index += 1
    return blocks


def xlsx_to_blocks(payload: bytes) -> list[CanonicalBlock]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("openpyxl is required to parse xlsx files") from exc

    workbook = load_workbook(filename=BytesIO(payload), data_only=True)
    blocks: list[CanonicalBlock] = []
    order_index = 0
    for sheet in workbook.worksheets:
        heading_path = [sheet.title]
        blocks.append(
            make_block(
                order_index=order_index,
                block_type="sheet",
                text=sheet.title,
                heading=sheet.title,
                heading_path=heading_path,
            )
        )
        order_index += 1
        for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
            values = [str(value).strip() for value in row if value not in (None, "")]
            if not values:
                continue
            blocks.append(
                make_block(
                    order_index=order_index,
                    block_type="table-row",
                    text=" | ".join(values),
                    heading=sheet.title,
                    heading_path=heading_path,
                    table_id=f"sheet-{sheet.title}",
                    row_index=row_index,
                    metadata={"table_columns": values, "sheet_name": sheet.title},
                )
            )
            order_index += 1
    return blocks


def pptx_to_blocks(payload: bytes) -> list[CanonicalBlock]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to parse pptx files") from exc

    presentation = Presentation(BytesIO(payload))
    blocks: list[CanonicalBlock] = []
    order_index = 0
    for slide_index, slide in enumerate(presentation.slides):
        slide_heading = f"Slide {slide_index + 1}"
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = " ".join(str(shape.text).split())
                if text:
                    slide_texts.append(text)
        if not slide_texts:
            continue
        for text in slide_texts:
            blocks.append(
                make_block(
                    order_index=order_index,
                    block_type="slide-text",
                    text=text,
                    heading=slide_heading,
                    heading_path=[slide_heading],
                    metadata={"slide_number": slide_index + 1},
                )
            )
            order_index += 1
    return blocks


def eml_to_blocks(payload: bytes) -> list[CanonicalBlock]:
    message = BytesParser(policy=policy.default).parsebytes(payload)
    blocks: list[CanonicalBlock] = []
    order_index = 0

    header_lines = []
    for header_name in ("subject", "from", "to", "date"):
        value = message.get(header_name)
        if value:
            header_lines.append(f"{header_name.title()}: {value}")
    if header_lines:
        blocks.append(
            make_block(
                order_index=order_index,
                block_type="email-header",
                text="\n".join(header_lines),
                heading="Email Header",
                heading_path=["Email Header"],
            )
        )
        order_index += 1

    body_parts: list[tuple[str, str]] = []
    if message.is_multipart():
        for part in message.walk():
            content_type = part.get_content_type()
            if content_type == "text/plain":
                body_parts.append(("email-body", part.get_content()))
            elif content_type == "text/html":
                body_parts.append(("email-html", html_to_text(part.get_content())))
    else:
        body_parts.append(("email-body", message.get_content()))

    for block_type, body_text in body_parts:
        for block in text_to_blocks(str(body_text), block_type):
            blocks.append(block.model_copy(update={"order_index": order_index}))
            order_index += 1
    return blocks


def make_block(
    *,
    order_index: int,
    block_type: str,
    text: str,
    heading: str | None = None,
    heading_path: list[str] | None = None,
    page_start: int | None = None,
    page_end: int | None = None,
    source_offset_start: int | None = None,
    source_offset_end: int | None = None,
    table_id: str | None = None,
    row_index: int | None = None,
    cell_index: int | None = None,
    parse_quality_score: float | None = None,
    ocr_confidence: float | None = None,
    metadata: dict | None = None,
) -> CanonicalBlock:
    return CanonicalBlock(
        id=str(uuid4()),
        block_type=block_type,
        order_index=order_index,
        heading=heading,
        heading_path=heading_path or [],
        text=text,
        page_start=page_start,
        page_end=page_end,
        source_offset_start=source_offset_start,
        source_offset_end=source_offset_end,
        table_id=table_id,
        row_index=row_index,
        cell_index=cell_index,
        parse_quality_score=parse_quality_score,
        ocr_confidence=ocr_confidence,
        metadata=metadata or {},
    )


def finalize_offsets(blocks: list[CanonicalBlock]) -> None:
    offset = 0
    for block in blocks:
        if block.source_offset_start is None:
            block.source_offset_start = offset
        if block.source_offset_end is None:
            block.source_offset_end = block.source_offset_start + len(block.text)
        offset = block.source_offset_end + 2


def calculate_parse_quality(
    *,
    blocks: list[CanonicalBlock],
    warnings: list[str],
    ocr_required: bool,
    ocr_applied: bool,
) -> float:
    if not blocks:
        return 0.0
    quality = 0.92
    if warnings:
        quality -= min(0.35, len(warnings) * 0.08)
    if ocr_required:
        quality -= 0.12
    if ocr_applied:
        quality -= 0.05
    average_block_quality = [
        block.parse_quality_score for block in blocks if block.parse_quality_score is not None
    ]
    if average_block_quality:
        quality = min(quality, sum(average_block_quality) / len(average_block_quality))
    return max(0.1, min(1.0, quality))


def parse_heading_level(style_name: str) -> int:
    match = re.search(r"(\d+)", style_name)
    if not match:
        return 1
    return max(1, int(match.group(1)))


def build_canonical_document(
    *,
    document_id: str,
    document_version_id: str,
    tenant_id: str,
    title: str,
    file_name: str,
    content_type: str,
    payload: bytes,
) -> CanonicalDocument:
    blocks, ocr_required, ocr_applied, parse_metadata = extract_text_blocks(
        file_name, content_type, payload
    )
    plain_text = "\n\n".join(block.text for block in blocks).strip()
    if not plain_text and ocr_required:
        blocks = [
            make_block(
                order_index=0,
                block_type="ocr-placeholder",
                text="OCR was required but no text could be extracted.",
                parse_quality_score=0.1,
                metadata={"ocr_status": "failed_or_unavailable"},
            )
        ]
        plain_text = blocks[0].text

    parse_warnings = [str(item) for item in parse_metadata.get("warnings", [])]
    parse_quality_score = float(parse_metadata.get("parse_quality_score", 0.0))
    return CanonicalDocument(
        document_id=document_id,
        document_version_id=document_version_id,
        tenant_id=tenant_id,
        title=title,
        source_format=file_name.rsplit(".", 1)[-1].lower() if "." in file_name else content_type,
        language=detect_language(plain_text),
        ocr_required=ocr_required,
        ocr_applied=ocr_applied,
        parse_quality_score=parse_quality_score,
        parse_warnings=parse_warnings,
        plain_text=plain_text,
        metadata={
            "file_name": file_name,
            "content_type": content_type,
            "block_count": len(blocks),
            "parser_format": str(parse_metadata.get("format", "")),
        },
        blocks=blocks,
    )
